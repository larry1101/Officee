import os
from io import BytesIO

from docx import Document
from docx.document import Document as TYPE_DOCUMENT
from docx.shared import Inches
from docx.styles.style import BaseStyle
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as TYPE_PRESENTATION


def convert_ppt(ppt_file_path, dest_file_path, printable=False):
    import win32com
    from win32com.client import Dispatch
    ppt = win32com.client.DispatchEx('PowerPoint.Application')
    ppt.Visible = True
    pptSel = ppt.Presentations.Open(ppt_file_path)
    win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
    docx = Document()
    for idx, slide in enumerate(pptSel.Slides):
        shapes = sorted(slide.Shapes, key=lambda s: s.Top)
        for shape in shapes:
            if shape.HasTextFrame:

                txt = shape.TextFrame.TextRange.Text
                # print('txt',txt,'b',txt.encode())
                if len(txt) == 0:
                    continue

                # 存在非ascii或unicode字符
                assert isinstance(txt, str)
                txt = txt.replace('\x0b', '[#UnASCIIorUnicode]')

                _ = docx.add_paragraph(txt)
                _.style.font.name = 'Arial'
                # 中文无法应用字体
                if printable:
                    print('=', txt, '@', _.style.name, '#', _.style.font.name)
            else:
                if printable:
                    print('not text frame')
                shape.Export(
                    r'E:\Python Projects\Officee\ppt\tmp.jpg',
                    1  # i is jpeg
                )
                docx.add_picture(
                    'tmp.jpg',
                    width=Inches(shape.Width/128)
                )
                if printable:
                    print('Image')
                    print()
                    # todo OCR
    pptSel.Close()
    ppt.Quit()
    docx.save(dest_file_path)


def convert(pptx, dest_file_name, printable=False):
    assert isinstance(pptx, TYPE_PRESENTATION)
    assert isinstance(dest_file_name, str)
    docx = Document()
    i = 0

    for slide in pptx.slides:
        if printable:
            print('\nNew Slide--------------', i + 1)
        # sort shapes by top
        sorted_shapes = sorted(list(slide.shapes), key=lambda x: x.top)
        # for shape in slide.shapes:
        for shape in sorted_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with open('tmp.jpg', 'wb') as f:
                    f.write(shape.image.blob)
                docx.add_picture(
                    'tmp.jpg',
                    # width=Inches(5.7)
                    width=Inches(shape.image.size[0]/256)
                )
                if printable:
                    print('Image')
                    print()
                # todo OCR
                continue
            if not shape.has_text_frame:
                continue
            pars = ''
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    pars += run.text

            pars = [p.split('  ') for p in pars.split('\u3000')]
            # print(pars)
            for par in pars:
                for seg in par:
                    if not seg == '' or not seg.replace(' ', '') == '':
                        _ = docx.add_paragraph(seg)
                        _.style.font.name = 'Arial'
                        # 中文无法应用字体
                        if printable:
                            print('=', seg, '@', _.style.name, '#', _.style.font.name)

        i += 1
        # if i>24:
        #     break

    docx.save(dest_file_name)

